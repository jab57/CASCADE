"""
Generate an impressive customer-facing PowerPoint presentation for GREmLN MCP Server
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with widescreen aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Modern Professional Color Scheme - "Midnight Science"
DARK_BLUE = RGBColor(0x1a, 0x1a, 0x2e)      # Deep midnight blue
MID_BLUE = RGBColor(0x16, 0x21, 0x3e)       # Navy
LIGHT_BLUE = RGBColor(0x4a, 0x69, 0xbd)     # Soft blue
ACCENT_TEAL = RGBColor(0x00, 0xd4, 0xaa)    # Vibrant teal/mint
ACCENT_GREEN = RGBColor(0x00, 0xc9, 0x8d)   # Emerald
ACCENT_ORANGE = RGBColor(0xff, 0x6b, 0x6b)  # Coral red
ACCENT_PURPLE = RGBColor(0x9d, 0x4e, 0xdd)  # Purple accent
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf8, 0xf9, 0xfc)     # Almost white
MID_GRAY = RGBColor(0x6c, 0x75, 0x7d)       # Neutral gray
DARK_GRAY = RGBColor(0x2c, 0x3e, 0x50)      # Charcoal


def add_gradient_background(slide, prs):
    """Add a modern subtle background with accent."""
    # Main background - clean white
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.fill.background()

    # Left accent strip - modern gradient effect
    left_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), prs.slide_height)
    left_strip.fill.solid()
    left_strip.fill.fore_color.rgb = ACCENT_TEAL
    left_strip.line.fill.background()

    # Bottom right decorative element
    corner = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(11.5), Inches(6.3), Inches(1.833), Inches(1.2))
    corner.fill.solid()
    corner.fill.fore_color.rgb = RGBColor(0xe8, 0xf4, 0xf8)
    corner.line.fill.background()
    corner.rotation = 180


def add_title_slide(prs, title, subtitle, tagline=None):
    """Add an impressive title slide with modern design."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Dark gradient background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Modern geometric accents - layered shapes
    accent1 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(-3), Inches(5), Inches(10), Inches(3))
    accent1.fill.solid()
    accent1.fill.fore_color.rgb = MID_BLUE
    accent1.line.fill.background()

    accent2 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(7), Inches(-1), Inches(8), Inches(2.5))
    accent2.fill.solid()
    accent2.fill.fore_color.rgb = ACCENT_TEAL
    accent2.line.fill.background()

    # Subtle circle decoration
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(5), Inches(3), Inches(3))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = ACCENT_PURPLE
    circle1.fill.fore_color.brightness = 0.3
    circle1.line.fill.background()

    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(-1), Inches(2.5), Inches(2.5))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = ACCENT_TEAL
    circle2.fill.fore_color.brightness = 0.4
    circle2.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.733), Inches(1.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.733), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_TEAL
    p.alignment = PP_ALIGN.LEFT

    # Tagline
    if tagline:
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.733), Inches(0.6))
        tf = tag_box.text_frame
        p = tf.paragraphs[0]
        p.text = tagline
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = RGBColor(0xaa, 0xbb, 0xcc)
        p.alignment = PP_ALIGN.LEFT

    return slide


def add_section_slide(prs, title, subtitle=None):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = MID_BLUE
    bg.line.fill.background()

    # Large accent shape
    accent = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3), Inches(-3), Inches(10), Inches(10))
    accent.fill.solid()
    accent.fill.fore_color.rgb = DARK_BLUE
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(4), Inches(2.8), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(4), Inches(4.3), Inches(9), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(22)
        p.font.color.rgb = ACCENT_TEAL
        p.alignment = PP_ALIGN.LEFT

    return slide


def add_content_slide(prs, title, bullets):
    """Add a slide with title and bullet points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_gradient_background(slide, prs)

    # Header area
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.4))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Bullets with icons
    start_y = 1.8
    for i, bullet in enumerate(bullets):
        y = Inches(start_y + i * 1.0)

        # Bullet icon (circle)
        icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.15), Inches(0.25), Inches(0.25))
        icon.fill.solid()
        icon.fill.fore_color.rgb = ACCENT_TEAL
        icon.line.fill.background()

        # Text
        text_box = slide.shapes.add_textbox(Inches(1.1), y, Inches(11.433), Inches(0.9))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = bullet
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_GRAY

    return slide


def add_icon_box(slide, x, y, width, height, icon_text, label, color):
    """Add a box with icon and label."""
    # Box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    box.adjustments[0] = 0.1

    # Icon text (emoji or symbol)
    icon_box = slide.shapes.add_textbox(x, y + Inches(0.15), width, Inches(0.6))
    tf = icon_box.text_frame
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Label
    label_box = slide.shapes.add_textbox(x, y + Inches(0.65), width, Inches(0.5))
    tf = label_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_stat_box(slide, x, y, number, label, color):
    """Add a statistics highlight box."""
    # Number
    num_box = slide.shapes.add_textbox(x, y, Inches(2.5), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    # Label
    label_box = slide.shapes.add_textbox(x, y + Inches(0.9), Inches(2.5), Inches(0.6))
    tf = label_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_ALIGN.CENTER


def add_arrow_shape(slide, x, y, direction="down"):
    """Add a directional arrow."""
    if direction == "down":
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, Inches(0.4), Inches(0.35))
    elif direction == "right":
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, Inches(0.5), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = MID_GRAY
    arrow.line.fill.background()
    return arrow


def add_flow_box(slide, x, y, width, height, text, fill_color, text_color=WHITE, font_size=14):
    """Add a rounded rectangle with centered text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.15

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    return shape


# ============ SLIDE 1: TITLE ============
add_title_slide(
    prs,
    "GREmLN",
    "AI-Powered Gene Perturbation Analysis",
    "Accelerating Drug Discovery with In Silico Predictions"
)

# ============ SLIDE 2: THE CHALLENGE ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, prs)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.4))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE
header.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.9))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "The Challenge"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# Challenge boxes
challenges = [
    ("$2.6B", "Average cost to\nbring a drug to market"),
    ("90%", "Failure rate in\nclinical trials"),
    ("10+ Years", "Typical drug\ndevelopment timeline"),
    ("1000s", "Genes to screen\nper indication")
]

for i, (stat, desc) in enumerate(challenges):
    x = Inches(0.8 + i * 3.1)

    # Box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(2.8), Inches(2.4))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.color.rgb = RGBColor(0xdd, 0xdd, 0xdd)
    box.line.width = Pt(1)

    # Stat
    stat_box = slide.shapes.add_textbox(x, Inches(2.3), Inches(2.8), Inches(1))
    tf = stat_box.text_frame
    p = tf.paragraphs[0]
    p.text = stat
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Description
    desc_box = slide.shapes.add_textbox(x, Inches(3.3), Inches(2.8), Inches(1))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

# Bottom message
msg_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(12.133), Inches(1))
tf = msg_box.text_frame
p = tf.paragraphs[0]
p.text = "What if you could predict perturbation effects before going to the lab?"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = MID_BLUE
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 3: THE SOLUTION ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, prs)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.4))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE
header.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.9))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "The Solution: GREmLN"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# Three pillars
pillars = [
    (ACCENT_TEAL, "Regulatory\nNetworks", "Pre-computed gene\nregulatory networks from\nsingle-cell data"),
    (MID_BLUE, "AI\nEmbeddings", "256-dim gene vectors\ntrained on 11 million\nsingle cells"),
    (ACCENT_GREEN, "External\nDatabases", "STRING protein interactions\nLINCS perturbation data\nSuper-enhancer annotations")
]

for i, (color, title, desc) in enumerate(pillars):
    x = Inches(0.9 + i * 4.1)

    # Pillar box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.8), Inches(3.8))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    # Title
    t_box = slide.shapes.add_textbox(x, Inches(2.1), Inches(3.8), Inches(1))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Description
    d_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(3.3), Inches(3.4), Inches(2))
    tf = d_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# Bottom tagline
tag_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.0), Inches(12.133), Inches(0.8))
tf = tag_box.text_frame
p = tf.paragraphs[0]
p.text = "Combined scoring delivers predictions no single approach can match"
p.font.size = Pt(22)
p.font.italic = True
p.font.color.rgb = MID_BLUE
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 4: HOW IT WORKS - ARCHITECTURE ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, prs)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.4))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE
header.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.9))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "How It Works"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# INPUT ROW
add_flow_box(slide, Inches(4.9), Inches(1.6), Inches(3.5), Inches(0.7),
             "Natural Language Query", DARK_GRAY, WHITE, 16)

# Example query
q_box = slide.shapes.add_textbox(Inches(4.9), Inches(2.35), Inches(3.5), Inches(0.5))
tf = q_box.text_frame
p = tf.paragraphs[0]
p.text = '"Knock down MYC in T cells"'
p.font.size = Pt(12)
p.font.italic = True
p.font.color.rgb = MID_GRAY
p.alignment = PP_ALIGN.CENTER

add_arrow_shape(slide, Inches(6.45), Inches(2.85), "down")

# PROCESSING ROW - Three parallel paths
add_flow_box(slide, Inches(0.5), Inches(3.3), Inches(3.5), Inches(1.1),
             "Network Propagation\nBFS through regulatory edges", MID_BLUE, WHITE, 14)
add_flow_box(slide, Inches(4.9), Inches(3.3), Inches(3.5), Inches(1.1),
             "Embedding Analysis\nCosine similarity in latent space", ACCENT_TEAL, WHITE, 14)
add_flow_box(slide, Inches(9.3), Inches(3.3), Inches(3.5), Inches(1.1),
             "External Validation\nSTRING + LINCS lookup", ACCENT_GREEN, WHITE, 14)

# Arrows down
add_arrow_shape(slide, Inches(2.05), Inches(4.5), "down")
add_arrow_shape(slide, Inches(6.45), Inches(4.5), "down")
add_arrow_shape(slide, Inches(10.85), Inches(4.5), "down")

# COMBINATION ROW
add_flow_box(slide, Inches(2.5), Inches(4.95), Inches(8.3), Inches(0.8),
             "Combined Scoring: Network Effect + Embedding Boost + Confidence", MID_BLUE, WHITE, 16)

add_arrow_shape(slide, Inches(6.45), Inches(5.85), "down")

# OUTPUT ROW
add_flow_box(slide, Inches(3.4), Inches(6.3), Inches(6.5), Inches(0.9),
             "Ranked Gene List + Actionable Recommendations", ACCENT_ORANGE, WHITE, 18)

# ============ SLIDE 5: PERTURBATION ALGORITHM - VERTICAL FLOWCHART ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, prs)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE
header.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(12.333), Inches(0.6))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Perturbation Analysis: Cascading Fallback"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE

# Flowchart Colors - Modern palette
DIAMOND_COLOR = RGBColor(0x2c, 0x3e, 0x50)  # Charcoal
YES_COLOR = ACCENT_TEAL
NO_COLOR = RGBColor(0x8e, 0x99, 0xa4)
END_COLOR = ACCENT_GREEN
BFS_COLOR = LIGHT_BLUE

# Layout constants
diamond_x = Inches(0.8)
diamond_w = Inches(2.0)
diamond_h = Inches(0.7)
bfs_x = Inches(4.5)
end_x = Inches(7.2)
row_spacing = Inches(0.85)
start_y = Inches(1.05)

# Helper to add a decision row
def add_decision_row(slide, y, step_num, question, action_text):
    # Decision diamond
    diamond = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, diamond_x, y, diamond_w, diamond_h)
    diamond.fill.solid()
    diamond.fill.fore_color.rgb = DIAMOND_COLOR
    diamond.line.color.rgb = WHITE
    diamond.line.width = Pt(1)

    # Question text (to the right of diamond, inside it won't fit well)
    q_box = slide.shapes.add_textbox(diamond_x + Inches(0.15), y + Inches(0.18), diamond_w - Inches(0.3), diamond_h - Inches(0.2))
    tf = q_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = question
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Step number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, diamond_x - Inches(0.35), y + Inches(0.2), Inches(0.32), Inches(0.32))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_TEAL
    badge.line.fill.background()
    badge_text = slide.shapes.add_textbox(diamond_x - Inches(0.35), y + Inches(0.22), Inches(0.32), Inches(0.3))
    tf = badge_text.text_frame
    p = tf.paragraphs[0]
    p.text = str(step_num)
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # YES arrow pointing right
    yes_arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, diamond_x + diamond_w + Inches(0.1), y + Inches(0.25), Inches(0.5), Inches(0.22))
    yes_arrow.fill.solid()
    yes_arrow.fill.fore_color.rgb = YES_COLOR
    yes_arrow.line.fill.background()

    # YES label
    yes_label = slide.shapes.add_textbox(diamond_x + diamond_w + Inches(0.15), y + Inches(0.02), Inches(0.4), Inches(0.25))
    tf = yes_label.text_frame
    p = tf.paragraphs[0]
    p.text = "YES"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = YES_COLOR

    # Action box (Route/BFS)
    action_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bfs_x, y + Inches(0.1), Inches(2.5), Inches(0.5))
    action_box.fill.solid()
    action_box.fill.fore_color.rgb = BFS_COLOR
    action_box.line.fill.background()

    action_label = slide.shapes.add_textbox(bfs_x + Inches(0.1), y + Inches(0.18), Inches(2.3), Inches(0.4))
    tf = action_label.text_frame
    p = tf.paragraphs[0]
    p.text = action_text
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Arrow to END
    end_arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, bfs_x + Inches(2.6), y + Inches(0.25), Inches(0.4), Inches(0.2))
    end_arrow.fill.solid()
    end_arrow.fill.fore_color.rgb = END_COLOR
    end_arrow.line.fill.background()

    # END box
    end_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, end_x, y + Inches(0.12), Inches(0.7), Inches(0.45))
    end_box.fill.solid()
    end_box.fill.fore_color.rgb = END_COLOR
    end_box.line.fill.background()

    end_label = slide.shapes.add_textbox(end_x, y + Inches(0.2), Inches(0.7), Inches(0.35))
    tf = end_label.text_frame
    p = tf.paragraphs[0]
    p.text = "END"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # NO arrow pointing down (except for last row)
    return y + diamond_h

# Gene Input box at top
input_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, diamond_x + Inches(0.3), start_y, Inches(1.4), Inches(0.45))
input_box.fill.solid()
input_box.fill.fore_color.rgb = ACCENT_ORANGE
input_box.line.fill.background()

input_label = slide.shapes.add_textbox(diamond_x + Inches(0.3), start_y + Inches(0.08), Inches(1.4), Inches(0.35))
tf = input_label.text_frame
p = tf.paragraphs[0]
p.text = "Gene Input"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Down arrow from input
down_arrow1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, diamond_x + Inches(0.85), start_y + Inches(0.48), Inches(0.3), Inches(0.3))
down_arrow1.fill.solid()
down_arrow1.fill.fore_color.rgb = NO_COLOR
down_arrow1.line.fill.background()

# Decision rows
y = start_y + Inches(0.85)

# Row 1: GRN edges
add_decision_row(slide, y, 1, "GRN edges?", "BFS propagation")
# NO arrow down
no_arrow1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, diamond_x + Inches(0.85), y + diamond_h + Inches(0.02), Inches(0.3), Inches(0.25))
no_arrow1.fill.solid()
no_arrow1.fill.fore_color.rgb = NO_COLOR
no_arrow1.line.fill.background()
no_label1 = slide.shapes.add_textbox(diamond_x + Inches(1.2), y + diamond_h + Inches(0.02), Inches(0.3), Inches(0.2))
tf = no_label1.text_frame
p = tf.paragraphs[0]
p.text = "NO"
p.font.size = Pt(7)
p.font.color.rgb = NO_COLOR

# Row 2: STRING
y += row_spacing + Inches(0.1)
add_decision_row(slide, y, 2, "STRING\npartner?", "Route → BFS")
no_arrow2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, diamond_x + Inches(0.85), y + diamond_h + Inches(0.02), Inches(0.3), Inches(0.25))
no_arrow2.fill.solid()
no_arrow2.fill.fore_color.rgb = NO_COLOR
no_arrow2.line.fill.background()
no_label2 = slide.shapes.add_textbox(diamond_x + Inches(1.2), y + diamond_h + Inches(0.02), Inches(0.3), Inches(0.2))
tf = no_label2.text_frame
p = tf.paragraphs[0]
p.text = "NO"
p.font.size = Pt(7)
p.font.color.rgb = NO_COLOR

# Row 3: Embedding neighbor
y += row_spacing + Inches(0.1)
add_decision_row(slide, y, 3, "Embedding\nneighbor?", "Route → BFS")
no_arrow3 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, diamond_x + Inches(0.85), y + diamond_h + Inches(0.02), Inches(0.3), Inches(0.25))
no_arrow3.fill.solid()
no_arrow3.fill.fore_color.rgb = NO_COLOR
no_arrow3.line.fill.background()
no_label3 = slide.shapes.add_textbox(diamond_x + Inches(1.2), y + diamond_h + Inches(0.02), Inches(0.3), Inches(0.2))
tf = no_label3.text_frame
p = tf.paragraphs[0]
p.text = "NO"
p.font.size = Pt(7)
p.font.color.rgb = NO_COLOR

# Row 4: LINCS
y += row_spacing + Inches(0.1)
add_decision_row(slide, y, 4, "LINCS\nsignature?", "Map downstream → BFS")
no_arrow4 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, diamond_x + Inches(0.85), y + diamond_h + Inches(0.02), Inches(0.3), Inches(0.25))
no_arrow4.fill.solid()
no_arrow4.fill.fore_color.rgb = NO_COLOR
no_arrow4.line.fill.background()
no_label4 = slide.shapes.add_textbox(diamond_x + Inches(1.2), y + diamond_h + Inches(0.02), Inches(0.3), Inches(0.2))
tf = no_label4.text_frame
p = tf.paragraphs[0]
p.text = "NO"
p.font.size = Pt(7)
p.font.color.rgb = NO_COLOR

# Row 5: Super-enhancer
y += row_spacing + Inches(0.1)
add_decision_row(slide, y, 5, "Super-\nenhancer?", "Route to TFs → BFS")
no_arrow5 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, diamond_x + Inches(0.85), y + diamond_h + Inches(0.02), Inches(0.3), Inches(0.25))
no_arrow5.fill.solid()
no_arrow5.fill.fore_color.rgb = NO_COLOR
no_arrow5.line.fill.background()
no_label5 = slide.shapes.add_textbox(diamond_x + Inches(1.2), y + diamond_h + Inches(0.02), Inches(0.3), Inches(0.2))
tf = no_label5.text_frame
p = tf.paragraphs[0]
p.text = "NO"
p.font.size = Pt(7)
p.font.color.rgb = NO_COLOR

# Row 6: Final fallback - Embedding only (no diamond, just result)
y += row_spacing + Inches(0.15)

# Step 6 badge
badge6 = slide.shapes.add_shape(MSO_SHAPE.OVAL, diamond_x - Inches(0.35), y + Inches(0.05), Inches(0.32), Inches(0.32))
badge6.fill.solid()
badge6.fill.fore_color.rgb = RGBColor(0x95, 0xa5, 0xa6)
badge6.line.fill.background()
badge6_text = slide.shapes.add_textbox(diamond_x - Inches(0.35), y + Inches(0.07), Inches(0.32), Inches(0.3))
tf = badge6_text.text_frame
p = tf.paragraphs[0]
p.text = "6"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Final result box
final_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, diamond_x, y, Inches(4.5), Inches(0.5))
final_box.fill.solid()
final_box.fill.fore_color.rgb = RGBColor(0x7f, 0x8c, 0x8d)
final_box.line.fill.background()

final_label = slide.shapes.add_textbox(diamond_x + Inches(0.1), y + Inches(0.1), Inches(4.3), Inches(0.4))
tf = final_label.text_frame
p = tf.paragraphs[0]
p.text = "Embedding-only similarity ranking"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Arrow to END
final_arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, diamond_x + Inches(4.6), y + Inches(0.15), Inches(0.4), Inches(0.2))
final_arrow.fill.solid()
final_arrow.fill.fore_color.rgb = RGBColor(0x95, 0xa5, 0xa6)
final_arrow.line.fill.background()

# Final END box (different color - lowest confidence)
final_end = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, diamond_x + Inches(5.1), y + Inches(0.02), Inches(1.2), Inches(0.45))
final_end.fill.solid()
final_end.fill.fore_color.rgb = RGBColor(0x95, 0xa5, 0xa6)
final_end.line.fill.background()

final_end_label = slide.shapes.add_textbox(diamond_x + Inches(5.1), y + Inches(0.1), Inches(1.2), Inches(0.35))
tf = final_end_label.text_frame
p = tf.paragraphs[0]
p.text = "END"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Right side: Legend/explanation box
legend_x = Inches(8.3)
legend_y = Inches(1.2)

legend_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, legend_x, legend_y, Inches(4.5), Inches(1.9))
legend_box.fill.solid()
legend_box.fill.fore_color.rgb = RGBColor(0xfa, 0xfa, 0xfa)
legend_box.line.color.rgb = LIGHT_GRAY
legend_box.line.width = Pt(1)

legend_title = slide.shapes.add_textbox(legend_x + Inches(0.2), legend_y + Inches(0.15), Inches(4.1), Inches(0.4))
tf = legend_title.text_frame
p = tf.paragraphs[0]
p.text = "Fallback Strategy"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

legend_content = slide.shapes.add_textbox(legend_x + Inches(0.2), legend_y + Inches(0.45), Inches(4.1), Inches(1.4))
tf = legend_content.text_frame
tf.word_wrap = True

steps_info = [
    ("1. GRN Edges", "Direct regulatory targets", ACCENT_TEAL),
    ("2. STRING", "Protein interaction partners", LIGHT_BLUE),
    ("3. Embeddings", "Functionally similar genes", ACCENT_PURPLE),
    ("4. LINCS", "Perturbation signatures", ACCENT_ORANGE),
    ("5. Super-enhancer", "Chromatin-linked TFs", ACCENT_GREEN),
    ("6. Embed-only", "Lowest confidence fallback", MID_GRAY),
]

for i, (title, desc, color) in enumerate(steps_info):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"{title}: {desc}"
    p.font.size = Pt(9)
    p.font.bold = False
    p.font.color.rgb = color
    p.space_before = Pt(2) if i > 0 else Pt(0)

# Equation box at bottom of legend - Full perturbation formulas
eq_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, legend_x + Inches(0.1), legend_y + Inches(3.1), Inches(4.3), Inches(3.5))
eq_box.fill.solid()
eq_box.fill.fore_color.rgb = DARK_BLUE
eq_box.line.fill.background()

eq_title = slide.shapes.add_textbox(legend_x + Inches(0.2), legend_y + Inches(3.18), Inches(4.1), Inches(0.3))
tf = eq_title.text_frame
p = tf.paragraphs[0]
p.text = "Perturbation Scoring"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL

# Network Effect equation
eq_net_label = slide.shapes.add_textbox(legend_x + Inches(0.2), legend_y + Inches(3.45), Inches(4.1), Inches(0.25))
tf = eq_net_label.text_frame
p = tf.paragraphs[0]
p.text = "Network Effect (BFS propagation):"
p.font.size = Pt(9)
p.font.color.rgb = RGBColor(0x88, 0xcc, 0xbb)

eq_net = slide.shapes.add_textbox(legend_x + Inches(0.2), legend_y + Inches(3.65), Inches(4.1), Inches(0.35))
tf = eq_net.text_frame
p = tf.paragraphs[0]
p.text = "E_net(t) = \u220f MI(e)  for edges r\u2192t"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Consolas"

# Embedding Effect equation
eq_emb_label = slide.shapes.add_textbox(legend_x + Inches(0.2), legend_y + Inches(4.0), Inches(4.1), Inches(0.25))
tf = eq_emb_label.text_frame
p = tf.paragraphs[0]
p.text = "Embedding Similarity:"
p.font.size = Pt(9)
p.font.color.rgb = RGBColor(0xbb, 0x88, 0xdd)

eq_emb = slide.shapes.add_textbox(legend_x + Inches(0.2), legend_y + Inches(4.2), Inches(4.1), Inches(0.35))
tf = eq_emb.text_frame
p = tf.paragraphs[0]
p.text = "E_emb(t) = cos(v_g, v_t)"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Consolas"

# Combined Effect equation
eq_comb_label = slide.shapes.add_textbox(legend_x + Inches(0.2), legend_y + Inches(4.55), Inches(4.1), Inches(0.25))
tf = eq_comb_label.text_frame
p = tf.paragraphs[0]
p.text = "Combined Effect:"
p.font.size = Pt(9)
p.font.color.rgb = RGBColor(0xff, 0xcc, 0x66)

eq_comb = slide.shapes.add_textbox(legend_x + Inches(0.2), legend_y + Inches(4.75), Inches(4.1), Inches(0.5))
tf = eq_comb.text_frame
p = tf.paragraphs[0]
p.text = "E(t) = E_net \u00d7 (\u03b1 + (1-\u03b1) \u00d7 E_emb)"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Consolas"

# Alpha value and explanation
eq_alpha = slide.shapes.add_textbox(legend_x + Inches(0.2), legend_y + Inches(5.15), Inches(4.1), Inches(0.25))
tf = eq_alpha.text_frame
p = tf.paragraphs[0]
p.text = "where \u03b1 = 0.7 (network weight)"
p.font.size = Pt(9)
p.font.italic = True
p.font.color.rgb = RGBColor(0xaa, 0xbb, 0xcc)

# Variable definitions
eq_def = slide.shapes.add_textbox(legend_x + Inches(0.2), legend_y + Inches(5.45), Inches(4.1), Inches(1.1))
tf = eq_def.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "r = Route(g) \u2192 GRN-active node"
p.font.size = Pt(8)
p.font.color.rgb = ACCENT_TEAL
p = tf.add_paragraph()
p.text = "v_g, v_t = 256-dim gene embeddings"
p.font.size = Pt(8)
p.font.color.rgb = RGBColor(0x99, 0xaa, 0xbb)
p = tf.add_paragraph()
p.text = "MI(e) = mutual information edge weight"
p.font.size = Pt(8)
p.font.color.rgb = RGBColor(0x99, 0xaa, 0xbb)

# ============ SLIDE 6: CELL TYPES ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, prs)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.4))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE
header.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.9))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "10 Cell Type Networks"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# Cell type grid
cell_types = [
    ("CD4+ T Cells", ACCENT_TEAL),
    ("CD8+ T Cells", ACCENT_TEAL),
    ("NK Cells", ACCENT_TEAL),
    ("NKT Cells", ACCENT_TEAL),
    ("B Cells", MID_BLUE),
    ("CD14+ Monocytes", MID_BLUE),
    ("CD16+ Monocytes", MID_BLUE),
    ("Dendritic Cells", MID_BLUE),
    ("Epithelial", ACCENT_GREEN),
    ("Erythrocytes", ACCENT_GREEN),
]

for i, (name, color) in enumerate(cell_types):
    row = i // 5
    col = i % 5
    x = Inches(0.6 + col * 2.5)
    y = Inches(1.8 + row * 2.4)

    # Cell icon (hexagon-ish)
    cell = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(2.0))
    cell.fill.solid()
    cell.fill.fore_color.rgb = color
    cell.line.fill.background()

    # Name
    n_box = slide.shapes.add_textbox(x, y + Inches(0.7), Inches(2.2), Inches(0.8))
    tf = n_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# Legend
legend_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.6), Inches(12.133), Inches(0.5))
tf = legend_box.text_frame
p = tf.paragraphs[0]
p.text = "Pre-computed regulatory networks derived from single-cell RNA-seq data"
p.font.size = Pt(16)
p.font.italic = True
p.font.color.rgb = MID_GRAY
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 6: CAPABILITIES ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, prs)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.4))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE
header.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.9))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Capabilities"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# Capability cards - 2 rows of 3
capabilities = [
    ("Gene Knockdown", "Predict downstream effects\nof gene silencing", ACCENT_TEAL),
    ("Overexpression", "Model increased gene\nexpression impact", MID_BLUE),
    ("Similar Genes", "Find functionally related\ngenes via embeddings", ACCENT_GREEN),
    ("Drug Targets", "Identify network\nvulnerabilities", ACCENT_ORANGE),
    ("Protein Interactions", "STRING database\nintegration", MID_BLUE),
    ("Super-Enhancers", "BRD4/BET inhibitor\nsensitivity", ACCENT_TEAL),
]

for i, (title, desc, color) in enumerate(capabilities):
    row = i // 3
    col = i % 3
    x = Inches(0.6 + col * 4.2)
    y = Inches(1.7 + row * 2.7)

    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.9), Inches(2.4))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = color
    card.line.width = Pt(3)

    # Color bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.9), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    # Title
    t_box = slide.shapes.add_textbox(x, y + Inches(0.4), Inches(3.9), Inches(0.6))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Description
    d_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.1), Inches(3.5), Inches(1.2))
    tf = d_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(15)
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 7: EXAMPLE WORKFLOW ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, prs)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.4))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE
header.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.9))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Example: APC Knockdown Analysis"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# Step boxes
steps = [
    ("1", "Query", '"What happens when\nAPC is knocked down?"', DARK_GRAY),
    ("2", "Classify", "Gene type: Scaffold\nNo transcriptional targets", MID_BLUE),
    ("3", "Suggest", "Analyze protein partners:\nCTNNB1, AXIN1, GSK3B", ACCENT_TEAL),
    ("4", "Predict", "CTNNB1 overexpression\n2,739 genes affected", ACCENT_GREEN),
    ("5", "Insight", "APC loss = beta-catenin\naccumulation = oncogenesis", ACCENT_ORANGE),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.4 + i * 2.55)

    # Step box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(2.4), Inches(3.2))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), Inches(2.0), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = WHITE
    circle.line.fill.background()

    n_box = slide.shapes.add_textbox(x + Inches(0.9), Inches(2.05), Inches(0.6), Inches(0.55))
    tf = n_box.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    # Title
    t_box = slide.shapes.add_textbox(x, Inches(2.7), Inches(2.4), Inches(0.5))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Description
    d_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(3.3), Inches(2.2), Inches(1.5))
    tf = d_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Arrow between boxes
    if i < len(steps) - 1:
        add_arrow_shape(slide, x + Inches(2.45), Inches(3.3), "right")

# Bottom insight
insight_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.3), Inches(12.133), Inches(1.5))
tf = insight_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Intelligent routing guides users from scaffold proteins to transcriptionally active partners"
p.font.size = Pt(20)
p.font.italic = True
p.font.color.rgb = MID_BLUE
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 8: USE CASES ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, prs)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.4))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE
header.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.9))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Use Cases"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# Use case cards
use_cases = [
    ("Immunotherapy", "Discover targets in tumor\nmicroenvironment across\nimmune cell types", ACCENT_TEAL),
    ("CAR-T Engineering", "Understand exhaustion\nand persistence pathways\nin T cell networks", MID_BLUE),
    ("CRISPR Prioritization", "Predict perturbation\neffects before costly\nwet lab validation", ACCENT_GREEN),
    ("Target Validation", "Compare candidates\nusing network vulnerability\nand druggability scores", ACCENT_ORANGE),
]

for i, (title, desc, color) in enumerate(use_cases):
    x = Inches(0.5 + i * 3.2)

    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.7), Inches(3.0), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()

    # Title
    t_box = slide.shapes.add_textbox(x, Inches(2.0), Inches(3.0), Inches(0.8))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Separator line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.5), Inches(2.85), Inches(2.0), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = WHITE
    line.line.fill.background()

    # Description
    d_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(3.1), Inches(2.6), Inches(2.5))
    tf = d_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 9: PERFORMANCE ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, prs)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.4))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE
header.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.9))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Performance"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# Stats
stats = [
    ("< 3 sec", "Full Analysis", "GPU-accelerated\nperturbation prediction"),
    ("19,247", "Genes", "Complete coverage in\nembedding space"),
    ("11M", "Cells Trained", "Deep learning on\nmassive single-cell data"),
    ("10", "Cell Types", "Immune & cancer-relevant\nnetworks"),
]

for i, (num, label, desc) in enumerate(stats):
    x = Inches(0.6 + i * 3.2)

    # Stat box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(2.9), Inches(3.5))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.fill.background()

    # Number
    n_box = slide.shapes.add_textbox(x, Inches(2.0), Inches(2.9), Inches(1.2))
    tf = n_box.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.alignment = PP_ALIGN.CENTER

    # Label
    l_box = slide.shapes.add_textbox(x, Inches(3.2), Inches(2.9), Inches(0.5))
    tf = l_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Description
    d_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(3.8), Inches(2.5), Inches(1.2))
    tf = d_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_ALIGN.CENTER

# Bottom bar
bottom = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.6), Inches(12.133), Inches(1.0))
bottom.fill.solid()
bottom.fill.fore_color.rgb = MID_BLUE
bottom.line.fill.background()

b_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.85), Inches(12.133), Inches(0.6))
tf = b_box.text_frame
p = tf.paragraphs[0]
p.text = "Caching and lazy loading ensure fast repeated queries"
p.font.size = Pt(20)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 10: INTEGRATION ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, prs)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.4))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE
header.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.9))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Seamless Integration"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# Integration diagram
# User
add_flow_box(slide, Inches(0.5), Inches(2.5), Inches(2.5), Inches(1.2),
             "Researcher\n(Natural Language)", DARK_GRAY, WHITE, 16)

add_arrow_shape(slide, Inches(3.1), Inches(3.0), "right")

# Claude
add_flow_box(slide, Inches(3.7), Inches(2.5), Inches(2.5), Inches(1.2),
             "Claude AI\n(MCP Client)", MID_BLUE, WHITE, 16)

add_arrow_shape(slide, Inches(6.3), Inches(3.0), "right")

# GREmLN
add_flow_box(slide, Inches(6.9), Inches(2.5), Inches(2.5), Inches(1.2),
             "GREmLN\n(MCP Server)", ACCENT_TEAL, WHITE, 16)

add_arrow_shape(slide, Inches(9.5), Inches(3.0), "right")

# Results
add_flow_box(slide, Inches(10.1), Inches(2.5), Inches(2.5), Inches(1.2),
             "Predictions\n& Insights", ACCENT_GREEN, WHITE, 16)

# Features below
features = [
    "Model Context Protocol (MCP) standard",
    "Natural language interface - no coding required",
    "Extensible tool framework for custom analyses",
    "Claude Desktop and Claude Code compatible"
]

for i, feat in enumerate(features):
    y = Inches(4.3 + i * 0.7)

    # Check mark
    check = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.5), y, Inches(0.35), Inches(0.35))
    check.fill.solid()
    check.fill.fore_color.rgb = ACCENT_GREEN
    check.line.fill.background()

    # Feature text
    f_box = slide.shapes.add_textbox(Inches(3.0), y, Inches(8), Inches(0.5))
    tf = f_box.text_frame
    p = tf.paragraphs[0]
    p.text = feat
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_GRAY

# ============ SLIDE 11: SUMMARY ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, prs)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.4))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE
header.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.9))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Why GREmLN?"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# Value props
values = [
    ("Accelerate Discovery", "Predict perturbation effects in seconds, not weeks", ACCENT_TEAL),
    ("Reduce Costs", "Prioritize experiments before wet lab validation", MID_BLUE),
    ("Increase Confidence", "AI + network analysis for robust predictions", ACCENT_GREEN),
    ("Enable Innovation", "Natural language interface democratizes access", ACCENT_ORANGE),
]

for i, (title, desc, color) in enumerate(values):
    y = Inches(1.6 + i * 1.05)

    # Color bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(0.15), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    # Title
    t_box = slide.shapes.add_textbox(Inches(1.0), y, Inches(4), Inches(0.6))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color

    # Description
    d_box = slide.shapes.add_textbox(Inches(1.0), y + Inches(0.45), Inches(11), Inches(0.5))
    tf = d_box.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY

# CTA box
cta = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(5.8), Inches(7.333), Inches(1.0))
cta.fill.solid()
cta.fill.fore_color.rgb = ACCENT_TEAL
cta.line.fill.background()

cta_box = slide.shapes.add_textbox(Inches(3), Inches(6.05), Inches(7.333), Inches(0.6))
tf = cta_box.text_frame
p = tf.paragraphs[0]
p.text = "Ready to transform your drug discovery pipeline?"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 12: THANK YOU ============
add_title_slide(
    prs,
    "Thank You",
    "Let's Discuss Your Use Case",
    None
)

# Save
output_path = "GREmLN_MCP_Server_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
